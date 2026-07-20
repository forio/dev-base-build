#!/usr/bin/env python3
"""
Generates the styled PowerPoint *template* for the Bike Shop example.

The Epicenter PowerPoint service fills this template with data at download time.
It locates the things to fill by **shape name**, so the chart/table shape names
below MUST match the `name` fields sent in the DocumentShadow
(see src/query/powerpoint.ts):

    Table   'FinalStandings'
    Chart   'ProfitByParticipant'   (bar)
    Chart   'ProfitOverTime'        (line)
    Chart   'RevenueShare'          (pie)

Text placeholders use {{token}} syntax and are replaced from
`environment.parameters`:

    {{title}} {{subtitle}} {{groupName}} {{episodeLabel}}
    {{participantCount}} {{generatedOn}}
    {{topPerformer}} {{topProfit}} {{averageProfit}}

All slide DESIGN lives here; the placeholder data below is only so the template
looks presentable on its own.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_TICK_LABEL_POSITION,
)
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ── Palette (Bike Shop: deep navy + energetic orange) ──────────────────
NAVY = RGBColor(0x0F, 0x2A, 0x43)
NAVY_SOFT = RGBColor(0x1B, 0x3A, 0x5B)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
PINK = RGBColor(0xDB, 0x27, 0x77)
PURPLE = RGBColor(0x93, 0x33, 0xEA)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF1, 0x5F, 0x5F)  # unused reserve
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLOUD = RGBColor(0xF3, 0xF6, 0xFB)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
SERIES = [BLUE, ORANGE, GREEN, PINK, PURPLE]

TITLE_FONT = "Century Schoolbook"
BODY_FONT = "Calibri"

EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color):
    r = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    solid(r, color)
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_para(p, text, size, color, bold=False, font=BODY_FONT,
             align=PP_ALIGN.LEFT, italic=False, spacing=None):
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return r


def name_shape(shape, name):
    shape.name = name


# ══════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)

# Motif: three overlapping accent circles (wheels), top-right, subtle.
for i, (cx, col, alpha) in enumerate([(11.3, ORANGE, None), (12.2, BLUE, None)]):
    c = s.shapes.add_shape(9, Inches(cx), Inches(-0.9), Inches(2.4), Inches(2.4))
    solid(c, col)
    c.shadow.inherit = False

eyebrow_tb, eyebrow = textbox(s, 0.9, 2.15, 11, 0.5)
set_para(eyebrow.paragraphs[0], "EPICENTER SIMULATION DEBRIEF", 14, ORANGE,
         bold=True, font=BODY_FONT)

title_tb, title = textbox(s, 0.9, 2.6, 11.5, 1.6)
set_para(title.paragraphs[0], "{{title}}", 52, WHITE, bold=True, font=TITLE_FONT)

sub_tb, sub = textbox(s, 0.9, 4.15, 11.5, 0.8)
set_para(sub.paragraphs[0], "{{subtitle}}", 24, RGBColor(0xCA, 0xDC, 0xFC),
         font=BODY_FONT)

meta_tb, meta = textbox(s, 0.9, 5.5, 11.5, 0.6)
set_para(meta.paragraphs[0],
         "{{groupName}}   •   {{episodeLabel}}   •   {{participantCount}} players",
         16, RGBColor(0x9F, 0xB3, 0xCC), font=BODY_FONT)

foot_tb, foot = textbox(s, 0.9, 6.75, 11.5, 0.4)
set_para(foot.paragraphs[0], "Generated {{generatedOn}}", 11, MUTED, font=BODY_FONT)


# ══════════════════════════════════════════════════════════════════════
# Slide 2 — How to read this deck
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)

_, h2 = textbox(s, 0.9, 0.7, 11.5, 1.0)
set_para(h2.paragraphs[0], "How to Read This Deck", 40, NAVY, bold=True,
         font=TITLE_FONT)

_, lead = textbox(s, 0.9, 1.75, 11.5, 0.7)
set_para(lead.paragraphs[0],
         "Each participant ran a bike shop and set a price for several years. "
         "The three result slides that follow summarize how they did.",
         16, MUTED, font=BODY_FONT)

cards = [
    ("1", "Final Standings",
     "A ranked table of every player by total profit, with average price and "
     "total revenue for the rounds they played.", ORANGE),
    ("2", "Profit by Player",
     "A bar chart comparing each player's cumulative profit — the quickest read "
     "on who ran the strongest shop.", BLUE),
    ("3", "The Market",
     "Profit over time as a line per player, plus a pie showing each player's "
     "share of total revenue.", GREEN),
]
cw, gap = 3.7, 0.42
x0 = 0.9
for i, (num, head, body, col) in enumerate(cards):
    x = x0 + i * (cw + gap)
    card = s.shapes.add_shape(1, Inches(x), Inches(2.9), Inches(cw), Inches(3.4))
    solid(card, CLOUD)
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    card.shadow.inherit = False

    badge = s.shapes.add_shape(9, Inches(x + 0.35), Inches(3.25),
                               Inches(0.8), Inches(0.8))
    solid(badge, col)
    badge.shadow.inherit = False
    bt = badge.text_frame
    bt.word_wrap = True
    set_para(bt.paragraphs[0], num, 24, WHITE, bold=True, align=PP_ALIGN.CENTER,
             font=TITLE_FONT)

    _, ch = textbox(s, x + 0.35, 4.35, cw - 0.7, 0.7)
    set_para(ch.paragraphs[0], head, 20, NAVY, bold=True, font=TITLE_FONT)
    _, cb = textbox(s, x + 0.35, 5.05, cw - 0.7, 1.6)
    set_para(cb.paragraphs[0], body, 14, INK, font=BODY_FONT, spacing=1.15)

_, foot2 = textbox(s, 0.9, 6.95, 11.5, 0.35)
set_para(foot2.paragraphs[0], "Generated {{generatedOn}}", 11, MUTED,
         font=BODY_FONT)


# ══════════════════════════════════════════════════════════════════════
# Slide 3 — Final Standings (KPIs + named TABLE)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)

_, h3 = textbox(s, 0.9, 0.6, 11.5, 0.9)
set_para(h3.paragraphs[0], "Final Standings", 40, NAVY, bold=True,
         font=TITLE_FONT)

kpis = [("WINNER", "{{topPerformer}}", ORANGE),
        ("WINNING PROFIT", "{{topProfit}}", BLUE),
        ("AVERAGE PROFIT", "{{averageProfit}}", GREEN)]
kw, kgap = 3.7, 0.42
for i, (lab, val, col) in enumerate(kpis):
    x = 0.9 + i * (kw + kgap)
    card = s.shapes.add_shape(1, Inches(x), Inches(1.6), Inches(kw), Inches(1.35))
    solid(card, CLOUD)
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    card.shadow.inherit = False
    _, kl = textbox(s, x + 0.35, 1.78, kw - 0.7, 0.4)
    set_para(kl.paragraphs[0], lab, 12, col, bold=True, font=BODY_FONT)
    _, kv = textbox(s, x + 0.35, 2.12, kw - 0.7, 0.75)
    set_para(kv.paragraphs[0], val, 20, INK, bold=True, font=TITLE_FONT)

# Named table — placeholder rows the service overwrites.
headers = ["Rank", "Player", "Rounds", "Avg. Price", "Total Revenue", "Total Profit"]
rows = [
    ["1", "Sample Player A", "5", "$142.00", "$1,240,000", "$486,000"],
    ["2", "Sample Player B", "5", "$128.50", "$1,180,000", "$421,500"],
    ["3", "Sample Player C", "4", "$155.00", "$980,000", "$362,000"],
    ["4", "Sample Player D", "5", "$119.00", "$1,050,000", "$318,750"],
]
nrows, ncols = len(rows) + 1, len(headers)
gf = s.shapes.add_table(nrows, ncols, Inches(0.9), Inches(3.35),
                        Inches(11.53), Inches(3.2))
name_shape(gf, "FinalStandings")
table = gf.table
table.first_row = True
table.horz_banding = True
widths = [1.1, 4.0, 1.3, 1.7, 1.9, 1.9]
for c, w in enumerate(widths):
    table.columns[c].width = Inches(w)
for c, htext in enumerate(headers):
    cell = table.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    set_para(p, htext, 13, WHITE, bold=True, font=BODY_FONT,
             align=PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER)
for r, row in enumerate(rows, start=1):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if r % 2 else CLOUD
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        set_para(p, val, 12, INK, font=BODY_FONT,
                 align=PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# Slide 4 — Profit by Player (named BAR chart)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)

_, h4 = textbox(s, 0.9, 0.6, 11.5, 0.9)
set_para(h4.paragraphs[0], "Profit by Player", 40, NAVY, bold=True,
         font=TITLE_FONT)
_, h4s = textbox(s, 0.9, 1.5, 11.5, 0.5)
set_para(h4s.paragraphs[0], "Cumulative profit across every round played.",
         16, MUTED, font=BODY_FONT)

cd = CategoryChartData()
cd.categories = ["Player A", "Player B", "Player C", "Player D"]
cd.add_series("Total Profit", (486000, 421500, 362000, 318750))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(2.2),
                        Inches(11.5), Inches(4.7), cd)
name_shape(gf, "ProfitByParticipant")
chart = gf.chart
chart.has_legend = False
chart.has_title = False
plot = chart.plots[0]
plot.gap_width = 80
plot.vary_by_categories = True
for i, pt in enumerate(plot.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = SERIES[i % len(SERIES)]
chart.value_axis.has_major_gridlines = True
chart.value_axis.tick_labels.font.size = Pt(11)
chart.value_axis.tick_labels.number_format = '"$"#,##0'
chart.value_axis.tick_labels.number_format_is_linked = False
chart.category_axis.tick_labels.font.size = Pt(12)


# ══════════════════════════════════════════════════════════════════════
# Slide 5 — The Market (named LINE + PIE charts)
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)

_, h5 = textbox(s, 0.9, 0.6, 11.5, 0.9)
set_para(h5.paragraphs[0], "The Market", 40, NAVY, bold=True, font=TITLE_FONT)

_, l5a = textbox(s, 0.9, 1.55, 5.8, 0.5)
set_para(l5a.paragraphs[0], "Profit over time", 18, NAVY, bold=True,
         font=TITLE_FONT)
_, l5b = textbox(s, 6.9, 1.55, 5.5, 0.5)
set_para(l5b.paragraphs[0], "Revenue share", 18, NAVY, bold=True,
         font=TITLE_FONT)

# Line chart
ld = CategoryChartData()
ld.categories = ["Yr 1", "Yr 2", "Yr 3", "Yr 4", "Yr 5"]
ld.add_series("Player A", (60000, 130000, 240000, 360000, 486000))
ld.add_series("Player B", (55000, 120000, 210000, 320000, 421500))
ld.add_series("Player C", (48000, 110000, 200000, 300000, 362000))
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.9), Inches(2.15),
                        Inches(5.9), Inches(4.7), ld)
name_shape(gf, "ProfitOverTime")
lchart = gf.chart
lchart.has_title = False
lchart.has_legend = True
lchart.legend.position = XL_LEGEND_POSITION.BOTTOM
lchart.legend.include_in_layout = False
lchart.legend.font.size = Pt(10)
for i, srs in enumerate(lchart.plots[0].series):
    srs.format.line.color.rgb = SERIES[i % len(SERIES)]
    srs.format.line.width = Pt(2.25)
lchart.value_axis.tick_labels.font.size = Pt(10)
lchart.value_axis.tick_labels.number_format = '"$"#,##0'
lchart.value_axis.tick_labels.number_format_is_linked = False
lchart.category_axis.tick_labels.font.size = Pt(11)
# Keep the year labels pinned to the bottom even when profit goes negative,
# instead of riding along the x-axis where it crosses zero.
lchart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

# Pie chart
pd = CategoryChartData()
pd.categories = ["Player A", "Player B", "Player C", "Player D"]
pd.add_series("Total Revenue", (1240000, 1180000, 980000, 1050000))
gf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(7.0), Inches(2.15),
                        Inches(5.4), Inches(4.7), pd)
name_shape(gf, "RevenueShare")
pchart = gf.chart
pchart.has_title = False
pchart.has_legend = True
pchart.legend.position = XL_LEGEND_POSITION.RIGHT
pchart.legend.include_in_layout = False
pchart.legend.font.size = Pt(11)
pplot = pchart.plots[0]
pplot.has_data_labels = True
dl = pplot.data_labels
dl.show_percentage = True
dl.show_value = False
dl.show_category_name = False
dl.number_format = '0%'
dl.number_format_is_linked = False
dl.font.size = Pt(11)
dl.font.color.rgb = WHITE
for i, pt in enumerate(pplot.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = SERIES[i % len(SERIES)]


import sys
out = sys.argv[1] if len(sys.argv) > 1 else "results-template.pptx"
prs.save(out)
print("saved", out)
